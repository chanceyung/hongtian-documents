"""Unit tests for PPTX parser."""

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from app.parsers.pptx_parser import PptxParser
from app.models.unified_document import UnifiedDocument


class TestPptxParser:
    """Test suite for PPTX parser."""

    @pytest.fixture
    def sample_pptx_with_content(self, tmp_path: Path) -> Path:
        """Create a sample PPTX file with text box and table."""
        prs = Presentation()

        # Add a slide
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

        # Add text box
        text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        text_frame = text_box.text_frame
        text_frame.text = "Sample Title"
        p = text_frame.add_paragraph()
        p.text = "Sample paragraph content"
        p.font.size = Pt(18)

        # Add table
        table = slide.shapes.add_table(
            rows=3,
            cols=2,
            left=Inches(1),
            top=Inches(2.5),
            width=Inches(4),
            height=Inches(1.5)
        )
        # Note: table.columns is accessed via table.table.columns in newer python-pptx
        try:
            table.table.columns[0].width = Inches(2)
            table.table.columns[1].width = Inches(2)
        except AttributeError:
            # For older versions, try direct access
            pass

        # Populate table using the correct API
        tbl = table.table
        cell = tbl.cell(0, 0)
        cell.text = "Header 1"
        cell = tbl.cell(0, 1)
        cell.text = "Header 2"

        cell = tbl.cell(1, 0)
        cell.text = "Row 1 Col 1"
        cell = tbl.cell(1, 1)
        cell.text = "Row 1 Col 2"

        cell = tbl.cell(2, 0)
        cell.text = "Row 2 Col 1"
        cell = tbl.cell(2, 1)
        cell.text = "Row 2 Col 2"

        # Save file
        file_path = tmp_path / "sample.pptx"
        prs.save(file_path)
        return file_path

    @pytest.fixture
    def empty_pptx(self, tmp_path: Path) -> Path:
        """Create an empty PPTX file with blank slides."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide

        file_path = tmp_path / "empty.pptx"
        prs.save(file_path)
        return file_path

    @pytest.mark.asyncio
    async def test_parser_pptx_with_textbox_and_table_returns_unified_document(
        self,
        sample_pptx_with_content: Path
    ) -> None:
        """Test parsing PPTX with text box and table returns UnifiedDocument."""
        parser = PptxParser()
        result: UnifiedDocument = await parser.parse(sample_pptx_with_content, "session-test")

        # Verify type
        assert isinstance(result, UnifiedDocument)

        # Verify texts extracted
        assert len(result.texts) > 0, "Should extract text from text box"

        # Verify tables extracted
        assert len(result.tables) > 0, "Should extract table"

        # Verify parse method
        assert result.parse_method == "python-pptx", "Should use python-pptx parser"

    @pytest.mark.asyncio
    async def test_parser_pptx_textbox_content_correct(self, sample_pptx_with_content: Path) -> None:
        """Test that text box content is correctly extracted."""
        parser = PptxParser()
        result: UnifiedDocument = await parser.parse(sample_pptx_with_content, "session-test")

        # Check for expected text content
        text_content = " ".join([t.content for t in result.texts])
        assert "Sample Title" in text_content, "Should extract title"
        assert "Sample paragraph content" in text_content, "Should extract paragraph"

    @pytest.mark.asyncio
    async def test_parser_pptx_table_structure_correct(self, sample_pptx_with_content: Path) -> None:
        """Test that table structure is correctly extracted."""
        parser = PptxParser()
        result: UnifiedDocument = await parser.parse(sample_pptx_with_content, "session-test")

        assert len(result.tables) == 1, "Should extract exactly one table"

        table = result.tables[0]
        assert len(table.headers) == 2, "Table should have 2 headers"
        # TableElement uses data field (list of lists), not rows
        assert len(table.data) == 2, "Table should have 2 data rows (excluding header)"

        # Verify header content
        assert table.headers[0] == "Header 1", "Should extract first header"
        assert table.headers[1] == "Header 2", "Should extract second header"

        # Verify row content from data field
        assert table.data[0][0] == "Row 1 Col 1", "Should extract first row first cell"
        assert table.data[0][1] == "Row 1 Col 2", "Should extract first row second cell"

    @pytest.mark.asyncio
    async def test_parser_pptx_empty_document(self, empty_pptx: Path) -> None:
        """Test parsing empty PPTX (no shapes)."""
        parser = PptxParser()
        result: UnifiedDocument = await parser.parse(empty_pptx, "session-test")

        assert isinstance(result, UnifiedDocument)
        assert result.parse_method == "python-pptx"
        # Empty document should have no texts or tables
        assert len(result.texts) == 0, "Empty PPTX should have no texts"
        assert len(result.tables) == 0, "Empty PPTX should have no tables"

    @pytest.mark.asyncio
    async def test_parser_pptx_nonexistent_file(self, tmp_path: Path) -> None:
        """Test parsing non-existent PPTX file raises appropriate error."""
        parser = PptxParser()
        nonexistent_file = tmp_path / "nonexistent.pptx"

        # python-pptx raises a different error for nonexistent files
        with pytest.raises(Exception):
            await parser.parse(nonexistent_file, "session-test")

    @pytest.mark.asyncio
    async def test_parser_pptx_multiple_slides(self, tmp_path: Path) -> None:
        """Test parsing PPTX with multiple slides."""
        prs = Presentation()

        # First slide with text
        slide1 = prs.slides.add_slide(prs.slide_layouts[5])
        text_box1 = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        text_box1.text_frame.text = "First Slide"

        # Second slide with text
        slide2 = prs.slides.add_slide(prs.slide_layouts[5])
        text_box2 = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        text_box2.text_frame.text = "Second Slide"

        file_path = tmp_path / "multiple_slides.pptx"
        prs.save(file_path)

        parser = PptxParser()
        result: UnifiedDocument = await parser.parse(file_path, "session-test")

        # Should extract texts from both slides
        text_content = " ".join([t.content for t in result.texts])
        assert "First Slide" in text_content, "Should extract first slide text"
        assert "Second Slide" in text_content, "Should extract second slide text"
        assert result.total_pages == 2, "Should have 2 pages"

    @pytest.mark.asyncio
    async def test_parser_pptx_fingerprint_uniqueness(self, sample_pptx_with_content: Path) -> None:
        """Test that fingerprint is unique and consistent for same file."""
        parser = PptxParser()

        # Parse same file twice
        result1: UnifiedDocument = await parser.parse(sample_pptx_with_content, "session-test")
        result2: UnifiedDocument = await parser.parse(sample_pptx_with_content, "session-test")

        # Compute fingerprints using the method
        fp1 = result1.compute_fingerprint()
        fp2 = result2.compute_fingerprint()

        # Fingerprints should be identical for same file
        assert fp1.text_fingerprints == fp2.text_fingerprints, "Fingerprint should be consistent"

        # Fingerprint should not be empty
        assert len(fp1.text_fingerprints) > 0, "Fingerprint should not be empty"