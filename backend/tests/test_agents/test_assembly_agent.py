"""Tests for AssemblyAgent."""
from __future__ import annotations

import pytest
from pathlib import Path
from unittest.mock import AsyncMock, MagicMock, patch
from app.models import UnifiedDocument, MagazineEditPlan, EditAction, SlideEditPlan
from app.agents.assembly_agent import AssemblyAgent, AssemblyResult, GenerationMetadata


@pytest.fixture
def assembly_agent():
    """Create AssemblyAgent instance for testing."""
    with patch.object(AssemblyAgent, "_verify_logo_files"):
        return AssemblyAgent()


@pytest.fixture
def sample_doc():
    """Create sample UnifiedDocument for testing."""
    return UnifiedDocument(
        source_file="test.pptx",
        source_format="pptx",
        title="Test Document",
        texts=[
            {"id": "text-1", "content": "Sample text content", "page": 1},
            {"id": "text-2", "content": "More text content", "page": 2},
        ],
        images=[
            {
                "id": "img-1",
                "local_path": "/path/to/image.png",
                "page": 1,
                "width": 300,
                "height": 200,
                "hash": "abc123"
            }
        ],
        tables=[
            {
                "id": "table-1",
                "page": 1,
                "data": [["A", "B"], ["1", "2"]],
                "headers": ["Col1", "Col2"]
            }
        ],
        linkage=[
            {
                "text_id": "text-1",
                "asset_id": "img-1",
                "asset_type": "image",
                "strategy": "spatial",
                "confidence": 0.9
            }
        ],
        total_pages=2,
        parse_method="python-pptx",
        parse_warnings=[]
    )


@pytest.fixture
def sample_edit_plan():
    """Create sample MagazineEditPlan for testing."""
    return MagazineEditPlan(
        document_id="test-doc",
        template_id="modern",
        pages=[
            SlideEditPlan(
                page_number=1,
                template_page="cover",
                actions=[
                    EditAction(
                        type="replace_text",
                        target_selector="text-1",
                        source_id="text-1",
                        content="Sample text content"
                    ),
                    EditAction(
                        type="replace_image",
                        target_selector="img-1",
                        source_id="img-1",
                        content=""
                    )
                ]
            ),
            SlideEditPlan(
                page_number=2,
                template_page="text_only",
                actions=[
                    EditAction(
                        type="replace_text",
                        target_selector="text-2",
                        source_id="text-2",
                        content="More text content"
                    )
                ]
            )
        ],
        design_spec={
            "color_theme": {
                "background": "#1a1a2e"
            }
        },
        original_fingerprint={}
    )


@pytest.fixture
def temp_pdf_pages(tmp_path):
    """Create temporary PDF page files for testing."""
    from PyPDF2 import PdfWriter

    pages = []
    for i in range(2):
        page_path = tmp_path / f"page_{i}.pdf"
        writer = PdfWriter()
        writer.add_blank_page(width=595, height=842)
        writer.write(str(page_path))
        pages.append(page_path)
    return pages


@pytest.fixture
def temp_pptx_slides(tmp_path):
    """Create temporary PPTX slide files for testing."""
    from pptx import Presentation

    slides = []
    for i in range(2):
        slide_path = tmp_path / f"slide_{i}.pptx"
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        prs.save(str(slide_path))
        slides.append(slide_path)
    return slides


class TestLogoSelection:
    """Tests for _select_logo method."""

    def test_select_logo_dark_background(self, assembly_agent):
        """Test logo selection for dark background."""
        dark_colors = ["#1a1a2e", "#0f3460", "#16213e", "#0a0a0a", "#111111"]
        for color in dark_colors:
            logo_path = assembly_agent._select_logo(color)
            assert logo_path == assembly_agent.LOGO_WHITE

    def test_select_logo_light_background(self, assembly_agent):
        """Test logo selection for light background."""
        light_colors = ["#ffffff", "#f5f5f5", "#eeeeee", "#ffffff ", "# FFFFFF"]
        for color in light_colors:
            logo_path = assembly_agent._select_logo(color)
            assert logo_path == assembly_agent.LOGO_BLACK

    def test_select_logo_case_insensitive(self, assembly_agent):
        """Test that logo selection is case insensitive."""
        assert assembly_agent._select_logo("#FFFFFF") == assembly_agent.LOGO_BLACK
        assert assembly_agent._select_logo("#1a1a2E") == assembly_agent.LOGO_WHITE

    def test_select_logo_with_whitespace(self, assembly_agent):
        """Test that logo selection handles whitespace."""
        assert assembly_agent._select_logo(" #ffffff ") == assembly_agent.LOGO_BLACK
        assert assembly_agent._select_logo("#1a1a2e  ") == assembly_agent.LOGO_WHITE


class TestPdfMerge:
    """Tests for PDF merge operations."""

    @pytest.mark.asyncio
    async def test_assemble_pdf_multiple_pages(self, assembly_agent, sample_doc, sample_edit_plan, temp_pdf_pages, tmp_path):
        """Test assembling PDF with multiple pages."""
        output_path = tmp_path / "output.pdf"

        with patch.object(assembly_agent, "_embed_logo_pdf", return_value=True):
            with patch.object(assembly_agent, "_write_metadata_pdf", return_value=True):
                result = await assembly_agent.assemble_pdf(
                    temp_pdf_pages, sample_edit_plan, sample_doc, output_path, 5.0
                )

                assert isinstance(result, AssemblyResult)
                assert result.success
                assert result.format == "pdf"
                assert result.pages_merged == 2
                assert result.logo_embedded
                assert result.metadata_written
                assert output_path.exists()
                assert result.generation_info["document_id"] == "test-doc"
                assert result.generation_info["total_pages"] == 2

    @pytest.mark.asyncio
    async def test_assemble_pdf_empty_pages(self, assembly_agent, sample_doc, sample_edit_plan, tmp_path):
        """Test assembling PDF with no pages."""
        output_path = tmp_path / "output.pdf"

        with patch.object(assembly_agent, "_create_empty_pdf"):
            with patch.object(assembly_agent, "_embed_logo_pdf", return_value=False):
                with patch.object(assembly_agent, "_write_metadata_pdf", return_value=True):
                    result = await assembly_agent.assemble_pdf(
                        [], sample_edit_plan, sample_doc, output_path, 1.0
                    )

                    assert result.success
                    assert result.pages_merged == 1
                    assert len(result.warnings) > 0
                    assert "No rendered pages provided" in result.warnings[0]

    @pytest.mark.asyncio
    async def test_merge_pdf_pages(self, assembly_agent, temp_pdf_pages, tmp_path):
        """Test _merge_pdf_pages method."""
        output_path = tmp_path / "merged.pdf"

        assembly_agent._merge_pdf_pages(temp_pdf_pages, output_path)

        assert output_path.exists()

        # Verify merged PDF has correct number of pages
        from PyPDF2 import PdfReader
        reader = PdfReader(str(output_path))
        assert len(reader.pages) == 2

    @pytest.mark.asyncio
    async def test_create_empty_pdf(self, assembly_agent, tmp_path):
        """Test _create_empty_pdf method."""
        output_path = tmp_path / "empty.pdf"

        assembly_agent._create_empty_pdf(output_path)

        assert output_path.exists()

        from PyPDF2 import PdfReader
        reader = PdfReader(str(output_path))
        assert len(reader.pages) == 1
        # Verify page dimensions (A4)
        page = reader.pages[0]
        assert abs(page.mediabox.width - 595) < 10
        assert abs(page.mediabox.height - 842) < 10


class TestPptxMerge:
    """Tests for PPTX merge operations."""

    @pytest.mark.asyncio
    async def test_assemble_pptx_multiple_slides(self, assembly_agent, sample_doc, sample_edit_plan, temp_pptx_slides, tmp_path):
        """Test assembling PPTX with multiple slides."""
        output_path = tmp_path / "output.pptx"

        with patch.object(assembly_agent, "_embed_logo_pptx", return_value=True):
            with patch.object(assembly_agent, "_write_metadata_pptx", return_value=True):
                result = await assembly_agent.assemble_pptx(
                    temp_pptx_slides, sample_edit_plan, sample_doc, output_path, 5.0
                )

                assert isinstance(result, AssemblyResult)
                assert result.success
                assert result.format == "pptx"
                assert result.pages_merged == 2
                assert result.logo_embedded
                assert result.metadata_written
                assert output_path.exists()
                assert result.generation_info["document_id"] == "test-doc"
                assert result.generation_info["total_pages"] == 2

    @pytest.mark.asyncio
    async def test_assemble_pptx_single_slide(self, assembly_agent, sample_doc, sample_edit_plan, temp_pptx_slides, tmp_path):
        """Test assembling PPTX with single slide."""
        output_path = tmp_path / "output.pptx"

        with patch("shutil.copy2"):
            with patch.object(assembly_agent, "_embed_logo_pptx", return_value=True):
                with patch.object(assembly_agent, "_write_metadata_pptx", return_value=True):
                    result = await assembly_agent.assemble_pptx(
                        [temp_pptx_slides[0]], sample_edit_plan, sample_doc, output_path, 2.0
                    )

                    assert result.success
                    assert result.pages_merged == 1

    @pytest.mark.asyncio
    async def test_create_empty_pptx(self, assembly_agent, tmp_path):
        """Test _create_empty_pptx method."""
        output_path = tmp_path / "empty.pptx"

        assembly_agent._create_empty_pptx(output_path)

        assert output_path.exists()

        from pptx import Presentation
        prs = Presentation(str(output_path))
        assert len(prs.slides) == 1


class TestMetadata:
    """Tests for metadata generation and writing."""

    def test_build_generation_metadata(self, assembly_agent, sample_doc, sample_edit_plan, tmp_path):
        """Test _build_generation_metadata method."""
        output_path = tmp_path / "output.pdf"
        metadata = assembly_agent._build_generation_metadata(
            sample_doc, sample_edit_plan, output_path, 5.0, 2
        )

        assert isinstance(metadata, GenerationMetadata)
        assert metadata.document_id == "test-doc"
        assert metadata.source_file == "test.pptx"
        assert metadata.source_format == "pptx"
        assert metadata.template_id == "modern"
        assert metadata.total_pages == 2
        assert metadata.text_count == 2
        assert metadata.image_count == 1
        assert metadata.table_count == 1
        assert metadata.generation_time_seconds == 5.0

    def test_build_fidelity_report(self, assembly_agent, sample_doc, sample_edit_plan):
        """Test _build_fidelity_report method."""
        report = assembly_agent._build_fidelity_report(sample_doc, sample_edit_plan)

        assert "text_count" in report
        assert "image_count" in report
        assert "table_count" in report
        assert "total_chars" in report
        assert "text_coverage" in report
        assert "image_coverage" in report
        assert "linkage_count" in report
        assert report["text_count"] == 2
        assert report["image_count"] == 1
        assert report["table_count"] == 1
        assert report["linkage_count"] == 1

    @pytest.mark.asyncio
    async def test_write_metadata_pdf(self, assembly_agent, sample_doc, tmp_path):
        """Test _write_metadata_pdf method."""
        pdf_path = tmp_path / "test.pdf"

        # Create a simple PDF first
        from PyPDF2 import PdfWriter
        writer = PdfWriter()
        writer.add_blank_page(width=595, height=842)
        writer.write(str(pdf_path))

        metadata = GenerationMetadata(
            document_id="test-doc",
            source_file="test.pptx",
            source_format="pptx",
            template_id="modern",
            output_format="pdf",
            generated_at="2024-01-01T00:00:00",
            generation_time_seconds=5.0,
            total_pages=1,
            text_count=1,
            image_count=0,
            table_count=0,
        )

        result = assembly_agent._write_metadata_pdf(pdf_path, metadata)

        assert result
        assert pdf_path.exists()

        # Verify metadata was written
        from PyPDF2 import PdfReader
        reader = PdfReader(str(pdf_path))
        metadata_dict = reader.metadata
        assert metadata_dict.get("/Title") is not None
        assert metadata_dict.get("/Author") == "弘天文档 AI"

    @pytest.mark.asyncio
    async def test_write_metadata_pptx(self, assembly_agent, sample_doc, tmp_path):
        """Test _write_metadata_pptx method."""
        pptx_path = tmp_path / "test.pptx"

        # Create a simple PPTX first
        from pptx import Presentation
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        prs.save(str(pptx_path))

        metadata = GenerationMetadata(
            document_id="test-doc",
            source_file="test.pptx",
            source_format="pptx",
            template_id="modern",
            output_format="pptx",
            generated_at="2024-01-01T00:00:00",
            generation_time_seconds=5.0,
            total_pages=1,
            text_count=1,
            image_count=0,
            table_count=0,
        )

        result = assembly_agent._write_metadata_pptx(pptx_path, metadata)

        assert result
        assert pptx_path.exists()

        # Verify metadata was written
        prs = Presentation(str(pptx_path))
        assert prs.core_properties.title is not None
        assert prs.core_properties.author == "弘天文档 AI"
        assert prs.core_properties.comments is not None
        # Comments should be truncated to 255 chars and contain key info
        assert len(prs.core_properties.comments) <= 255
        assert "test-doc" in prs.core_properties.comments
        assert "modern" in prs.core_properties.comments


class TestResultModel:
    """Tests for AssemblyResult model."""

    def test_assembly_result_structure(self):
        """Test that AssemblyResult has correct structure."""
        result = AssemblyResult(
            success=True,
            output_path=Path("/output/test.pdf"),
            format="pdf",
            pages_merged=5,
            logo_embedded=True,
            metadata_written=True,
            generation_info={"test": "data"},
            fidelity_report={"score": 0.95},
            warnings=[]
        )

        assert result.success
        assert result.format == "pdf"
        assert result.pages_merged == 5
        assert result.logo_embedded
        assert result.metadata_written
        assert len(result.warnings) == 0

    def test_assembly_result_with_warnings(self):
        """Test AssemblyResult with warnings."""
        result = AssemblyResult(
            success=True,
            output_path=Path("/output/test.pdf"),
            format="pdf",
            pages_merged=5,
            logo_embedded=False,
            metadata_written=True,
            generation_info={},
            fidelity_report={},
            warnings=["Logo not found", "Metadata incomplete"]
        )

        assert result.success
        assert len(result.warnings) == 2
        assert "Logo not found" in result.warnings
        assert "Metadata incomplete" in result.warnings


class TestErrorHandling:
    """Tests for error handling."""

    @pytest.mark.asyncio
    async def test_assemble_pdf_with_missing_pages(self, assembly_agent, sample_doc, sample_edit_plan, tmp_path):
        """Test assembling PDF when some pages are missing."""
        output_path = tmp_path / "output.pdf"
        missing_pages = [tmp_path / "nonexistent.pdf"]

        result = await assembly_agent.assemble_pdf(
            missing_pages, sample_edit_plan, sample_doc, output_path, 1.0
        )

        # Should handle gracefully - might create empty PDF or merge available pages
        assert isinstance(result, AssemblyResult)
        assert output_path.exists()

    @pytest.mark.asyncio
    async def test_assembly_result_on_error(self, assembly_agent, sample_doc, sample_edit_plan, tmp_path):
        """Test that AssemblyResult contains error information when assembly fails."""
        output_path = tmp_path / "output.pdf"

        # Force an error by using invalid data
        with patch.object(assembly_agent, "_merge_pdf_pages", side_effect=Exception("Test error")):
            result = await assembly_agent.assemble_pdf(
                [tmp_path / "page.pdf"], sample_edit_plan, sample_doc, output_path, 1.0
            )

            assert isinstance(result, AssemblyResult)
            assert not result.success
            assert len(result.warnings) > 0
            assert "Assembly failed" in result.warnings[0]