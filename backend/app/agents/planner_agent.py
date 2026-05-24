"""Planner Agent — 文档复杂度评估与自适应路径选择"""
from __future__ import annotations

from pathlib import Path
from typing import Literal

from app.core.logging import get_logger
from app.models.execution_plan import ComplexityMetrics, ExecutionPlan

logger = get_logger(__name__)

PathType = Literal["fast", "standard", "deep"]


class PlannerAgent:
    """上传后自动评估复杂度，生成执行计划"""

    async def plan(self, file_path: Path) -> ExecutionPlan:
        metrics = await self._quick_scan(file_path)
        score = self._complexity_score(metrics)
        path = self._select_path(score)
        return self._generate_plan(metrics, score, path)

    async def _quick_scan(self, file_path: Path) -> ComplexityMetrics:
        ext = file_path.suffix.lower()
        if ext == ".pptx":
            return await self._scan_pptx(file_path)
        if ext == ".pdf":
            return await self._scan_pdf(file_path)
        if ext in (".docx", ".xlsx", ".md", ".txt"):
            return await self._scan_generic(file_path)
        return ComplexityMetrics()

    async def _scan_pptx(self, path: Path) -> ComplexityMetrics:
        from pptx import Presentation

        prs = Presentation(str(path))
        page_count = len(prs.slides)
        image_count = 0
        table_count = 0
        total_chars = 0
        layout_types: set[str] = set()

        for slide in prs.slides:
            has_image = False
            has_text = False
            for shape in slide.shapes:
                if shape.shape_type == 13:
                    image_count += 1
                    has_image = True
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        total_chars += len(para.text.strip())
                    if any(p.text.strip() for p in shape.text_frame.paragraphs):
                        has_text = True
                if shape.has_table:
                    table_count += 1
            if has_image and has_text:
                layout_types.add("text_image")
            elif has_image:
                layout_types.add("image_only")
            elif has_text:
                layout_types.add("text_only")

        text_density = total_chars / max(page_count, 1)
        image_ratio = image_count / max(page_count, 1)

        return ComplexityMetrics(
            page_count=page_count,
            image_count=image_count,
            table_count=table_count,
            total_chars=total_chars,
            text_density=text_density,
            image_ratio=image_ratio,
            layout_diversity=len(layout_types),
        )

    async def _scan_pdf(self, path: Path) -> ComplexityMetrics:
        import fitz

        doc = fitz.open(str(path))
        page_count = len(doc)
        image_count = 0
        total_chars = 0

        for page in doc:
            total_chars += len(page.get_text().strip())
            image_count += len(page.get_images(full=True))

        doc.close()

        text_density = total_chars / max(page_count, 1)
        image_ratio = image_count / max(page_count, 1)

        return ComplexityMetrics(
            page_count=page_count,
            image_count=image_count,
            table_count=0,
            total_chars=total_chars,
            text_density=text_density,
            image_ratio=image_ratio,
            layout_diversity=3,
        )

    async def _scan_generic(self, path: Path) -> ComplexityMetrics:
        size_mb = path.stat().st_size / (1024 * 1024)
        est_pages = max(1, int(size_mb * 2))

        return ComplexityMetrics(
            page_count=est_pages,
            image_count=0,
            table_count=0,
            total_chars=int(size_mb * 5000),
            text_density=500,
            image_ratio=0.0,
            layout_diversity=1,
        )

    def _complexity_score(self, m: ComplexityMetrics) -> float:
        score = 0.0

        score += min(m.page_count * 2, 30)
        score += min(m.image_count * 1.5, 25)
        score += min(m.table_count * 3, 15)
        score += min(m.text_density / 50, 15)
        score += min(m.image_ratio * 10, 10)
        score += min(m.layout_diversity * 3, 5)

        return min(score, 100.0)

    def _select_path(self, score: float) -> PathType:
        if score <= 30:
            return "fast"
        if score <= 70:
            return "standard"
        return "deep"

    def _generate_plan(
        self, metrics: ComplexityMetrics, score: float, path: PathType,
    ) -> ExecutionPlan:
        risks: list[str] = []

        if metrics.page_count > 30:
            risks.append(f"文档较大（{metrics.page_count}页），处理时间较长")
        if metrics.image_ratio > 2.0:
            risks.append("图片密度高，素材补充可能需要较长时间")
        if metrics.table_count > 5:
            risks.append("包含较多表格，渲染可能较慢")
        if metrics.text_density > 800:
            risks.append("文字密集，排版布局需仔细规划")

        if path == "fast":
            return ExecutionPlan(
                complexity_score=score,
                complexity_metrics=metrics,
                processing_path="fast",
                estimated_time_seconds=max(30, metrics.page_count * 5),
                estimated_api_calls=2,
                estimated_cost_cny=0.02,
                skip_analyzer=True,
                skip_supplement=True,
                page_parallel=False,
                max_render_concurrency=1,
                risk_alerts=risks,
            )

        if path == "deep":
            return ExecutionPlan(
                complexity_score=score,
                complexity_metrics=metrics,
                processing_path="deep",
                estimated_time_seconds=max(120, metrics.page_count * 15),
                estimated_api_calls=max(10, metrics.page_count),
                estimated_cost_cny=max(0.1, metrics.page_count * 0.02),
                skip_analyzer=False,
                skip_supplement=False,
                page_parallel=True,
                max_render_concurrency=4,
                risk_alerts=risks,
            )

        return ExecutionPlan(
            complexity_score=score,
            complexity_metrics=metrics,
            processing_path="standard",
            estimated_time_seconds=max(60, metrics.page_count * 10),
            estimated_api_calls=max(5, metrics.page_count // 3),
            estimated_cost_cny=max(0.05, metrics.page_count * 0.01),
            skip_analyzer=False,
            skip_supplement=False,
            page_parallel=False,
            max_render_concurrency=2,
            risk_alerts=risks,
        )
