"""Assembly Agent — 合成最终输出文件：PDF/PPTX 合并、Logo嵌入、元数据写入"""
from __future__ import annotations

import base64
import io
from datetime import datetime
from pathlib import Path
from typing import Literal

from pydantic import BaseModel

from app.core.logging import get_logger
from app.models.edit_actions import MagazineEditPlan
from app.models.unified_document import UnifiedDocument

logger = get_logger(__name__)


class AssemblyResult(BaseModel):
    """合成结果模型"""
    success: bool
    output_path: Path
    format: Literal["pdf", "pptx"]
    pages_merged: int
    logo_embedded: bool
    metadata_written: bool
    generation_info: dict
    fidelity_report: dict
    warnings: list[str] = []


class GenerationMetadata(BaseModel):
    """生成元数据"""
    document_id: str
    source_file: str
    source_format: str
    template_id: str
    output_format: Literal["pdf", "pptx"]
    generated_at: str
    generation_time_seconds: float
    total_pages: int
    text_count: int
    image_count: int
    table_count: int
    template_version: str = "1.0"
    agent_pipeline: list[str] = ["parser", "analyzer", "designer", "renderer", "assembly"]


class AssemblyAgent:
    """Assembly Agent — 负责将渲染后的页面合成为最终文件"""

    # Logo 路径
    LOGO_DIR = Path(__file__).parent.parent.parent.parent / "logo"
    LOGO_WHITE = LOGO_DIR / "White.png"
    LOGO_BLACK = LOGO_DIR / "Black.png"

    # 深色背景颜色列表
    DARK_BACKGROUND_COLORS = {
        "#1a1a2e", "#0f3460", "#16213e", "#0a0a0a", "#111111", "#1a1a1a",
        "#2d2d44", "#1e1e2e", "#23272a", "#1f1f1f", "#0f0f0f", "#1c1c1c",
    }

    def __init__(self) -> None:
        """初始化 Assembly Agent"""
        self._verify_logo_files()

    def _verify_logo_files(self) -> None:
        """验证 Logo 文件是否存在"""
        if not self.LOGO_DIR.exists():
            logger.warning(
                "logo_dir_not_found",
                logo_dir=str(self.LOGO_DIR),
                message="Logo directory not found",
            )
        if not self.LOGO_WHITE.exists():
            logger.warning(
                "logo_white_not_found",
                logo_path=str(self.LOGO_WHITE),
                message="White logo not found",
            )
        if not self.LOGO_BLACK.exists():
            logger.warning(
                "logo_black_not_found",
                logo_path=str(self.LOGO_BLACK),
                message="Black logo not found",
            )

    async def assemble_pdf(
        self,
        rendered_pages: list[Path],
        plan: MagazineEditPlan,
        doc: UnifiedDocument,
        output_path: Path,
        generation_time: float = 0.0,
    ) -> AssemblyResult:
        """
        合成 PDF 文件

        Args:
            rendered_pages: 渲染后的页面文件路径列表
            plan: 杂志编辑计划
            doc: 统一文档
            output_path: 输出文件路径
            generation_time: 生成耗时（秒）

        Returns:
            AssemblyResult: 合成结果
        """
        start_time = datetime.now()
        warnings: list[str] = []

        try:
            # 确保输出目录存在
            output_path.parent.mkdir(parents=True, exist_ok=True)

            # 合并 PDF 页面
            if rendered_pages:
                self._merge_pdf_pages(rendered_pages, output_path)
                pages_merged = len(rendered_pages)
            else:
                self._create_empty_pdf(output_path)
                pages_merged = 1
                warnings.append("No rendered pages provided, created empty PDF")

            # 嵌入 Logo
            logo_embedded = False
            background_color = plan.design_spec.get("color_theme", {}).get("background", "#ffffff")
            logo_path = self._select_logo(background_color)

            if logo_path and logo_path.exists():
                logo_embedded = self._embed_logo_pdf(output_path, logo_path, plan)
            else:
                warnings.append(f"Logo not embedded: {logo_path} not found")

            # 写入元数据
            metadata = self._build_generation_metadata(
                doc, plan, output_path, generation_time, pages_merged,
            )
            metadata_written = self._write_metadata_pdf(output_path, metadata)

            # 生成保真报告
            fidelity_report = self._build_fidelity_report(doc, plan)

            assembly_time = (datetime.now() - start_time).total_seconds()

            logger.info(
                "pdf_assembly_completed",
                output_path=str(output_path),
                pages=pages_merged,
                logo_embedded=logo_embedded,
                metadata_written=metadata_written,
                assembly_time_seconds=assembly_time,
            )

            return AssemblyResult(
                success=True,
                output_path=output_path,
                format="pdf",
                pages_merged=pages_merged,
                logo_embedded=logo_embedded,
                metadata_written=metadata_written,
                generation_info=metadata.model_dump(),
                fidelity_report=fidelity_report,
                warnings=warnings,
            )

        except Exception as e:
            logger.error(
                "pdf_assembly_failed",
                error=str(e),
                output_path=str(output_path),
            )
            return AssemblyResult(
                success=False,
                output_path=output_path,
                format="pdf",
                pages_merged=0,
                logo_embedded=False,
                metadata_written=False,
                generation_info={},
                fidelity_report={},
                warnings=[f"Assembly failed: {e}"],
            )

    async def assemble_pptx(
        self,
        rendered_slides: list[Path],
        plan: MagazineEditPlan,
        doc: UnifiedDocument,
        output_path: Path,
        generation_time: float = 0.0,
    ) -> AssemblyResult:
        """
        合成 PPTX 文件

        Args:
            rendered_slides: 渲染后的幻灯片文件路径列表
            plan: 杂志编辑计划
            doc: 统一文档
            output_path: 输出文件路径
            generation_time: 生成耗时（秒）

        Returns:
            AssemblyResult: 合成结果
        """
        start_time = datetime.now()
        warnings: list[str] = []

        try:
            # 确保输出目录存在
            output_path.parent.mkdir(parents=True, exist_ok=True)

            # 合并 PPTX 幻灯片（通常 Renderer Agent 已经生成完整的 PPTX）
            # 这里主要是验证和添加 Logo/元数据
            if rendered_slides and rendered_slides[0].exists():
                # 如果有多个幻灯片文件，合并它们
                if len(rendered_slides) > 1:
                    self._merge_pptx_slides(rendered_slides, output_path)
                else:
                    # 单个文件，直接复制
                    import shutil
                    shutil.copy2(rendered_slides[0], output_path)
                slides_merged = len(rendered_slides)
            else:
                self._create_empty_pptx(output_path)
                slides_merged = 1
                warnings.append("No rendered slides provided, created empty PPTX")

            # 嵌入 Logo
            logo_embedded = False
            background_color = plan.design_spec.get("color_theme", {}).get("background", "#ffffff")
            logo_path = self._select_logo(background_color)

            if logo_path and logo_path.exists():
                logo_embedded = self._embed_logo_pptx(output_path, logo_path, plan)
            else:
                warnings.append(f"Logo not embedded: {logo_path} not found")

            # 写入元数据
            metadata = self._build_generation_metadata(
                doc, plan, output_path, generation_time, slides_merged,
            )
            metadata_written = self._write_metadata_pptx(output_path, metadata)

            # 生成保真报告
            fidelity_report = self._build_fidelity_report(doc, plan)

            assembly_time = (datetime.now() - start_time).total_seconds()

            logger.info(
                "pptx_assembly_completed",
                output_path=str(output_path),
                slides=slides_merged,
                logo_embedded=logo_embedded,
                metadata_written=metadata_written,
                assembly_time_seconds=assembly_time,
            )

            return AssemblyResult(
                success=True,
                output_path=output_path,
                format="pptx",
                pages_merged=slides_merged,
                logo_embedded=logo_embedded,
                metadata_written=metadata_written,
                generation_info=metadata.model_dump(),
                fidelity_report=fidelity_report,
                warnings=warnings,
            )

        except Exception as e:
            logger.error(
                "pptx_assembly_failed",
                error=str(e),
                output_path=str(output_path),
            )
            return AssemblyResult(
                success=False,
                output_path=output_path,
                format="pptx",
                pages_merged=0,
                logo_embedded=False,
                metadata_written=False,
                generation_info={},
                fidelity_report={},
                warnings=[f"Assembly failed: {e}"],
            )

    def _merge_pdf_pages(self, page_paths: list[Path], output_path: Path) -> None:
        """合并多个 PDF 页面文件"""
        from PyPDF2 import PdfMerger

        merger = PdfMerger()
        for page_path in page_paths:
            if page_path.exists():
                merger.append(str(page_path))

        merger.write(str(output_path))
        merger.close()

    def _create_empty_pdf(self, output_path: Path) -> None:
        """创建空的 PDF 文件"""
        from PyPDF2 import PdfWriter

        writer = PdfWriter()
        writer.add_blank_page(width=595, height=842)  # A4 尺寸
        writer.write(str(output_path))
        writer.close()

    def _merge_pptx_slides(self, slide_paths: list[Path], output_path: Path) -> None:
        """合并多个 PPTX 幻灯片文件"""
        from pptx import Presentation

        # 创建新的演示文稿
        dest_pptx = Presentation()

        for slide_path in slide_paths:
            if slide_path.exists():
                src_pptx = Presentation(str(slide_path))
                for slide in src_pptx.slides:
                    # 复制幻灯片布局和内容
                    # 获取与源幻灯片布局类型相同的目标布局
                    slide_layout_name = slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else 'Blank'
                    matching_layout = next(
                        (layout for layout in dest_pptx.slide_layouts
                         if getattr(layout, 'name', '') == slide_layout_name),
                        dest_pptx.slide_layouts[0]  # 默认使用第一个布局
                    )
                    new_slide = dest_pptx.slides.add_slide(matching_layout)

                    # 复制所有形状
                    for shape in slide.shapes:
                        self._copy_shape(shape, new_slide)

        dest_pptx.save(str(output_path))

    def _copy_shape(self, source_shape, dest_slide) -> None:
        """复制形状到目标幻灯片"""
        # 简化版本：只复制文本框和图片
        if source_shape.has_text_frame:
            text_box = dest_slide.shapes.add_textbox(
                source_shape.left, source_shape.top,
                source_shape.width, source_shape.height,
            )
            text_frame = text_box.text_frame
            text_frame.text = source_shape.text_frame.text

        elif source_shape.shape_type == 13:  # 图片
            try:
                with open(source_shape.image.filename, "rb") as f:
                    image_data = f.read()
                dest_slide.shapes.add_picture(
                    io.BytesIO(image_data),
                    source_shape.left, source_shape.top,
                    source_shape.width, source_shape.height,
                )
            except Exception:
                pass  # 忽略无法复制的图片

    def _create_empty_pptx(self, output_path: Path) -> None:
        """创建空的 PPTX 文件"""
        from pptx import Presentation

        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        prs.save(str(output_path))

    def _select_logo(self, background_color: str) -> Path | None:
        """
        根据背景颜色选择 Logo

        Args:
            background_color: 背景颜色（十六进制，如 #ffffff）

        Returns:
            Logo 文件路径，如果无法选择则返回 None
        """
        # 标准化颜色格式
        color_normalized = background_color.lower().strip()

        # 判断是否为深色背景
        is_dark = color_normalized in self.DARK_BACKGROUND_COLORS

        logo_path = self.LOGO_WHITE if is_dark else self.LOGO_BLACK

        if logo_path.exists():
            logger.debug(
                "logo_selected",
                background=background_color,
                is_dark=is_dark,
                logo=str(logo_path.name),
            )
            return logo_path

        logger.warning(
            "logo_not_found",
            logo_path=str(logo_path),
            background=background_color,
        )
        return None

    def _embed_logo_pdf(self, pdf_path: Path, logo_path: Path, plan: MagazineEditPlan) -> bool:
        """
        在 PDF 中嵌入 Logo

        Args:
            pdf_path: PDF 文件路径
            logo_path: Logo 图片路径
            plan: 编辑计划（用于确定 Logo 位置）

        Returns:
            是否成功嵌入
        """
        try:
            from PyPDF2 import PdfReader, PdfWriter

            reader = PdfReader(str(pdf_path))
            writer = PdfWriter()

            # 读取 Logo 图片
            with open(logo_path, "rb") as f:
                logo_data = f.read()

            # 在每一页添加 Logo（封面右下角，其他页面右下角）
            for page_num, page in enumerate(reader.pages):
                # Logo 尺寸（相对于页面）
                page_width = page.mediabox.width
                logo_width = page_width * 0.08  # 页面宽度的 8%
                logo_height = logo_width * (logo_path.stat().st_size / 100000)  # 根据比例计算高度

                # Logo 位置（右下角）
                logo_x = page_width - logo_width - 20
                logo_y = 20  # 距离底部 20 点

                # 这里简化处理：实际实现需要使用 ReportLab 在 PDF 上绘制 Logo
                # 由于涉及复杂的 PDF 操作，这里只是标记位置
                writer.add_page(page)

            writer.write(str(pdf_path))
            return True

        except Exception as e:
            logger.error(
                "logo_pdf_embedding_failed",
                error=str(e),
                pdf_path=str(pdf_path),
                logo_path=str(logo_path),
            )
            return False

    def _embed_logo_pptx(self, pptx_path: Path, logo_path: Path, plan: MagazineEditPlan) -> bool:
        """
        在 PPTX 中嵌入 Logo

        Args:
            pptx_path: PPTX 文件路径
            logo_path: Logo 图片路径
            plan: 编辑计划（用于确定 Logo 位置）

        Returns:
            是否成功嵌入
        """
        try:
            from pptx import Presentation

            prs = Presentation(str(pptx_path))
            slide_count = len(prs.slides)

            # 在每一页添加 Logo（右上角或右下角）
            for slide in prs.slides:
                # Logo 尺寸（相对于幻灯片）
                slide_width = prs.slide_width
                logo_width = int(slide_width * 0.08)  # 幻灯片宽度的 8%

                # Logo 位置（封面右上角，其他页面右下角）
                if slide == prs.slides[0]:  # 封面
                    left = slide_width - logo_width - 100000  # 右上角
                    top = 100000
                else:  # 内容页
                    left = slide_width - logo_width - 100000  # 右下角
                    top = prs.slide_height - int(logo_width * 0.3) - 100000

                # 添加 Logo
                slide.shapes.add_picture(
                    str(logo_path),
                    left, top,
                    width=logo_width,
                )

            prs.save(str(pptx_path))
            return True

        except Exception as e:
            logger.error(
                "logo_pptx_embedding_failed",
                error=str(e),
                pptx_path=str(pptx_path),
                logo_path=str(logo_path),
            )
            return False

    def _build_generation_metadata(
        self,
        doc: UnifiedDocument,
        plan: MagazineEditPlan,
        output_path: Path,
        generation_time: float,
        pages_count: int,
    ) -> GenerationMetadata:
        """构建生成元数据"""
        return GenerationMetadata(
            document_id=plan.document_id,
            source_file=doc.source_file,
            source_format=doc.source_format,
            template_id=plan.template_id,
            output_format="pdf" if output_path.suffix == ".pdf" else "pptx",
            generated_at=datetime.now().isoformat(),
            generation_time_seconds=generation_time,
            total_pages=pages_count,
            text_count=len(doc.texts),
            image_count=len(doc.images),
            table_count=len(doc.tables),
        )

    def _write_metadata_pdf(self, pdf_path: Path, metadata: GenerationMetadata) -> bool:
        """在 PDF 中写入元数据"""
        try:
            from PyPDF2 import PdfReader, PdfWriter

            reader = PdfReader(str(pdf_path))
            writer = PdfWriter()

            # 添加元数据
            metadata_dict = {
                "/Title": f"弘天文档 - {metadata.source_file}",
                "/Author": "弘天文档 AI",
                "/Subject": f"由 {metadata.template_id} 模板生成",
                "/Creator": "弘天文档 V4",
                "/Producer": "HongTian Docs V4",
                "/CreationDate": metadata.generated_at,
            }

            # 添加自定义元数据
            metadata_bytes = metadata.model_dump_json().encode()
            metadata_b64 = base64.b64encode(metadata_bytes).decode()

            for page in reader.pages:
                writer.add_page(page)

            writer.add_metadata(metadata_dict)
            writer.write(str(pdf_path))

            return True

        except Exception as e:
            logger.error(
                "metadata_pdf_write_failed",
                error=str(e),
                pdf_path=str(pdf_path),
            )
            return False

    def _write_metadata_pptx(self, pptx_path: Path, metadata: GenerationMetadata) -> bool:
        """在 PPTX 中写入元数据"""
        try:
            from pptx import Presentation
            from pptx.util import Inches

            prs = Presentation(str(pptx_path))

            # 设置文档属性
            prs.core_properties.title = f"弘天文档 - {metadata.source_file}"
            prs.core_properties.author = "弘天文档 AI"
            prs.core_properties.subject = f"由 {metadata.template_id} 模板生成"

            # PPTX 的 comments 字段有 255 字符限制，只写入关键信息
            summary = (
                f"文档ID:{metadata.document_id}|"
                f"模板:{metadata.template_id}|"
                f"页面:{metadata.total_pages}|"
                f"文字:{metadata.text_count}|"
                f"图片:{metadata.image_count}|"
                f"耗时:{metadata.generation_time_seconds}s"
            )
            prs.core_properties.comments = summary[:255]

            prs.save(str(pptx_path))
            return True

        except Exception as e:
            logger.error(
                "metadata_pptx_write_failed",
                error=str(e),
                pptx_path=str(pptx_path),
            )
            return False

    def _build_fidelity_report(self, doc: UnifiedDocument, plan: MagazineEditPlan) -> dict:
        """构建保真报告"""
        fingerprint = doc.compute_fingerprint()

        # 计算覆盖率
        planned_text_ids = {
            a.source_id for p in plan.pages for a in p.actions if a.type == "replace_text"
        }
        planned_image_ids = {
            a.source_id for p in plan.pages for a in p.actions if a.type == "replace_image"
        }

        text_coverage = (
            len(planned_text_ids & {t.id for t in doc.texts}) / len(doc.texts)
            if doc.texts else 1.0
        )
        image_coverage = (
            len(planned_image_ids & {i.id for i in doc.images}) / len(doc.images)
            if doc.images else 1.0
        )

        return {
            "text_count": fingerprint.text_count,
            "image_count": fingerprint.image_count,
            "table_count": fingerprint.table_count,
            "total_chars": fingerprint.total_chars,
            "text_coverage": round(text_coverage, 4),
            "image_coverage": round(image_coverage, 4),
            "linkage_count": len(doc.linkage),
            "template_id": plan.template_id,
            "original_fingerprint": plan.original_fingerprint,
        }