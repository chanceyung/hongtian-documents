"""SVG → DrawingML → PPTX converter for PPT Master SVG templates."""

import base64
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Pt


@dataclass
class SvgElement:
    """Parsed SVG element with computed EMU coordinates."""
    tag: str
    attrs: dict[str, Any]
    x_emu: int
    y_emu: int
    width_emu: int
    height_emu: int


class SvgToPptxConverter:
    """Convert PPT Master SVG pages to PPTX presentation."""

    EMU_PER_PX = 12700

    def convert(
        self,
        svg_pages: list[str],
        design_spec: dict,
        output_path: Path,
        template_pptx: Path | None = None,
    ) -> Path:
        """Convert list of SVG pages to PPTX file.

        Args:
            svg_pages: List of SVG content strings
            design_spec: Design specification dict (for metadata)
            output_path: Output PPTX file path
            template_pptx: Optional template PPTX to use as base

        Returns:
            Path to generated PPTX file
        """
        if template_pptx and template_pptx.exists():
            prs = Presentation(str(template_pptx))
            prs.slides._sldIdLst.clear()
        else:
            prs = Presentation()
            prs.slide_width = Emu(12192000)
            prs.slide_height = Emu(6858000)

        for svg_content in svg_pages:
            self._process_svg_page(svg_content, prs, design_spec)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        return output_path

    def _process_svg_page(self, svg_content: str, prs: Presentation, design_spec: dict) -> None:
        """Process a single SVG page and add to presentation."""
        soup = BeautifulSoup(svg_content, "xml")
        svg = soup.find("svg")

        if not svg:
            return

        viewbox = svg.get("viewBox")
        if not viewbox:
            viewbox = svg.get("viewbox", "0 0 1920 1080")

        vb_parts = [float(x) for x in viewbox.split()]
        vb_x, vb_y, vb_w, vb_h = vb_parts

        slide = prs.slides.add_slide()
        slide_width_px = vb_w
        slide_height_px = vb_h

        slide.shapes._spTree.clear()

        for element in svg.find_all(recursive=True):
            if element.name == "svg":
                continue

            svg_elem = self._parse_svg_element(element, vb_w, vb_h)
            if svg_elem:
                self._add_element_to_slide(slide, svg_elem)

    def _parse_svg_element(self, element: Any, vb_w: float, vb_h: float) -> SvgElement | None:
        """Parse SVG element and compute EMU coordinates."""
        tag = element.name
        attrs = dict(element.attrs)

        x, y, width, height = self._extract_geometry(element, vb_w, vb_h)

        if width <= 0 or height <= 0:
            return None

        return SvgElement(
            tag=tag,
            attrs={**attrs, "_raw_element": element},
            x_emu=self._svg_to_emu(x, vb_w, int(12192000)),
            y_emu=self._svg_to_emu(y, vb_h, int(6858000)),
            width_emu=self._svg_to_emu(width, vb_w, int(12192000)),
            height_emu=self._svg_to_emu(height, vb_h, int(6858000)),
        )

    def _extract_geometry(self, element: Any, vb_w: float, vb_h: float) -> tuple[float, float, float, float]:
        """Extract x, y, width, height from element."""
        if element.name == "rect":
            return (
                float(element.get("x", 0)),
                float(element.get("y", 0)),
                float(element.get("width", vb_w)),
                float(element.get("height", vb_h)),
            )
        elif element.name == "image":
            return (
                float(element.get("x", 0)),
                float(element.get("y", 0)),
                float(element.get("width", vb_w)),
                float(element.get("height", vb_h)),
            )
        elif element.name in ("circle", "ellipse"):
            cx = float(element.get("cx", vb_w / 2))
            cy = float(element.get("cy", vb_h / 2))
            if element.name == "circle":
                r = float(element.get("r", min(vb_w, vb_h) / 4))
                return cx - r, cy - r, r * 2, r * 2
            else:
                rx = float(element.get("rx", vb_w / 4))
                ry = float(element.get("ry", vb_h / 4))
                return cx - rx, cy - ry, rx * 2, ry * 2
        elif element.name == "text":
            return (
                float(element.get("x", 0)),
                float(element.get("y", 0)),
                float(element.get("width", vb_w / 2)),
                40,
            )
        elif element.name == "line":
            x1 = float(element.get("x1", 0))
            y1 = float(element.get("y1", 0))
            x2 = float(element.get("x2", vb_w))
            y2 = float(element.get("y2", 0))
            return min(x1, x2), min(y1, y2), abs(x2 - x1), abs(y2 - y1)
        else:
            return 0, 0, vb_w, vb_h

    def _add_element_to_slide(self, slide: Any, element: SvgElement) -> None:
        """Add parsed element to PPTX slide."""
        if element.tag == "rect":
            self._add_rect(slide, element)
        elif element.tag == "text":
            self._add_text(slide, element)
        elif element.tag == "image":
            self._add_image(slide, element)
        elif element.tag in ("circle", "ellipse"):
            self._add_oval(slide, element)
        elif element.tag == "line":
            self._add_line(slide, element)

    def _add_rect(self, slide: Any, element: SvgElement) -> None:
        """Add rectangle shape to slide."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(element.x_emu),
            Emu(element.y_emu),
            Emu(element.width_emu),
            Emu(element.height_emu),
        )
        shape.line.color.rgb = self._parse_color_to_rgb(element.attrs.get("stroke", "none"))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._parse_color_to_rgb(element.attrs.get("fill", "none"))

        if "rx" in element.attrs or "ry" in element.attrs:
            rx = float(element.attrs.get("rx", 0))
            ry = float(element.attrs.get("ry", rx))
            shape.adjustments[0] = Emu(self._svg_to_emu(rx, 100, Emu(100000).emu))

    def _add_text(self, slide: Any, element: SvgElement) -> None:
        """Add textbox to slide."""
        shape = slide.shapes.add_textbox(
            Emu(element.x_emu),
            Emu(element.y_emu),
            Emu(element.width_emu),
            Emu(element.height_emu),
        )
        text_frame = shape.text_frame
        text_frame.word_wrap = True

        text_content = element.attrs.get("content", "")
        if not text_content:
            raw = element.attrs.get("_raw_element")
            if raw is not None:
                text_content = raw.get_text()

        p = text_frame.paragraphs[0]
        p.text = text_content or ""

        font_size = float(element.attrs.get("font-size", 12))
        p.font.size = Pt(font_size)

        color_str = element.attrs.get("fill", "#000000")
        p.font.color.rgb = self._parse_color_to_rgb(color_str)

        font_family = element.attrs.get("font-family", "Arial")
        safe_fonts = ["Arial", "Helvetica", "Calibri", "Times New Roman", "Verdana"]
        p.font.name = next((f for f in safe_fonts if f.lower() in font_family.lower()), "Arial")

        align_map = {"start": 0, "middle": 1, "end": 2}
        text_anchor = element.attrs.get("text-anchor", "start")
        p.alignment = align_map.get(text_anchor, 0)

    def _add_image(self, slide: Any, element: SvgElement) -> None:
        """Add image to slide."""
        href = element.attrs.get("href", "")
        if not href:
            href = element.attrs.get("{http://www.w3.org/1999/xlink}href", "")

        if href.startswith("data:"):
            img_data = self._decode_base64_image(href)
        else:
            return

        slide.shapes.add_picture(
            img_data,
            Emu(element.x_emu),
            Emu(element.y_emu),
            width=Emu(element.width_emu),
            height=Emu(element.height_emu),
        )

    def _decode_base64_image(self, data_uri: str) -> bytes:
        """Decode base64 data URI to image bytes."""
        match = re.match(r"data:image/([a-zA-Z]+);base64,(.+)", data_uri)
        if not match:
            raise ValueError(f"Invalid data URI: {data_uri[:50]}...")

        return base64.b64decode(match.group(2))

    def _add_oval(self, slide: Any, element: SvgElement) -> None:
        """Add oval shape to slide."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Emu(element.x_emu),
            Emu(element.y_emu),
            Emu(element.width_emu),
            Emu(element.height_emu),
        )
        shape.line.color.rgb = self._parse_color_to_rgb(element.attrs.get("stroke", "none"))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._parse_color_to_rgb(element.attrs.get("fill", "none"))

    def _add_line(self, slide: Any, element: SvgElement) -> None:
        """Add line to slide."""
        x1 = float(element.attrs.get("x1", 0))
        y1 = float(element.attrs.get("y1", 0))
        x2 = float(element.attrs.get("x2", element.width_emu))
        y2 = float(element.attrs.get("y2", element.height_emu))

        line = slide.shapes.add_connector(
            1,
            Emu(self._svg_to_emu(x1, 1920, 12192000)),
            Emu(self._svg_to_emu(y1, 1080, 6858000)),
            Emu(self._svg_to_emu(x2, 1920, 12192000)),
            Emu(self._svg_to_emu(y2, 1080, 6858000)),
        )
        line.line.color.rgb = self._parse_color_to_rgb(element.attrs.get("stroke", "#000000"))
        line.width = Pt(float(element.attrs.get("stroke-width", 1)))

    def _svg_to_emu(self, value: float, total_svg: float, total_emu: int) -> int:
        """Convert SVG coordinate to EMU."""
        return int((value / total_svg) * total_emu)

    def _parse_color(self, color_str: str) -> str:
        """Parse color string to hex format."""
        if not color_str or color_str in ("none", "transparent"):
            return "FFFFFF"

        color_str = color_str.strip()

        if color_str.startswith("#"):
            hex_str = color_str[1:]
            if len(hex_str) == 3:
                return f"{hex_str[0]}{hex_str[0]}{hex_str[1]}{hex_str[1]}{hex_str[2]}{hex_str[2]}"
            return hex_str.upper()

        rgb_match = re.match(r"rgb\((\d+),\s*(\d+),\s*(\d+)\)", color_str)
        if rgb_match:
            r, g, b = rgb_match.groups()
            return f"{int(r):02X}{int(g):02X}{int(b):02X}"

        rgba_match = re.match(r"rgba\((\d+),\s*(\d+),\s*(\d+),\s*[\d.]+\)", color_str)
        if rgba_match:
            r, g, b = rgba_match.groups()[:3]
            return f"{int(r):02X}{int(g):02X}{int(b):02X}"

        if color_str.startswith("url("):
            return "000000"

        color_map = {
            "black": "000000",
            "white": "FFFFFF",
            "red": "FF0000",
            "green": "00FF00",
            "blue": "0000FF",
            "yellow": "FFFF00",
            "cyan": "00FFFF",
            "magenta": "FF00FF",
        }
        return color_map.get(color_str.lower(), "000000")

    def _parse_color_to_rgb(self, color_str: str) -> RGBColor | None:
        """Parse color string to RGBColor object."""
        hex_color = self._parse_color(color_str)
        if hex_color == "FFFFFF" and color_str in ("none", "transparent"):
            return None

        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)
        return RGBColor(r, g, b)