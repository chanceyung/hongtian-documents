"""PDF/PPTX 生成引擎 - 本地运行"""
from fastapi import APIRouter, HTTPException
from pydantic import BaseModel
from pathlib import Path
import json

from app.core.config import settings

router = APIRouter(prefix="/generate", tags=["文档生成"])


class GenerateRequest(BaseModel):
    task_id: str
    session_id: str
    output_format: str = "pdf"  # pdf | pptx
    template_id: str = "default"


class GenerateResult(BaseModel):
    task_id: str
    output_path: str
    format: str
    pages: int


@router.post("/pdf", response_model=GenerateResult)
async def generate_pdf(req: GenerateRequest):
    """使用 Typst 生成杂志级 PDF"""
    # 加载排版指令
    layout_file = Path(settings.OUTPUT_DIR) / req.task_id / "layout.json"
    if not layout_file.exists():
        raise HTTPException(400, "排版指令不存在，请先完成排版规划")

    with open(layout_file, "r", encoding="utf-8") as f:
        layout = json.load(f)

    # 加载解析结果
    structured_file = Path(settings.OUTPUT_DIR) / req.task_id / "structured.json"
    with open(structured_file, "r", encoding="utf-8") as f:
        structured = json.load(f)

    # 生成 Typst 标记语言
    typst_content = _build_typst_document(layout, structured)

    # 编译 PDF
    output_path = Path(settings.OUTPUT_DIR) / req.task_id / "output.pdf"
    output_path.parent.mkdir(parents=True, exist_ok=True)

    typst_file = Path(settings.OUTPUT_DIR) / req.task_id / "document.typ"
    with open(typst_file, "w", encoding="utf-8") as f:
        f.write(typst_content)

    try:
        import subprocess
        result = subprocess.run(
            ["typst", "compile", str(typst_file), str(output_path)],
            capture_output=True,
            text=True,
            timeout=120,
        )
        if result.returncode != 0:
            raise HTTPException(500, f"PDF 编译失败: {result.stderr}")
    except FileNotFoundError:
        raise HTTPException(500, "Typst 未安装，请先安装: https://typst.app")

    return GenerateResult(
        task_id=req.task_id,
        output_path=str(output_path),
        format="pdf",
        pages=len(layout.get("pages", [])),
    )


@router.post("/pptx", response_model=GenerateResult)
async def generate_pptx(req: GenerateRequest):
    """使用 python-pptx 生成可编辑 PPTX"""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    layout_file = Path(settings.OUTPUT_DIR) / req.task_id / "layout.json"
    if not layout_file.exists():
        raise HTTPException(400, "排版指令不存在")

    with open(layout_file, "r", encoding="utf-8") as f:
        layout = json.load(f)

    structured_file = Path(settings.OUTPUT_DIR) / req.task_id / "structured.json"
    with open(structured_file, "r", encoding="utf-8") as f:
        structured = json.load(f)

    prs = Presentation()
    # 16:9 比例
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 构建素材索引
    image_map = {img["id"]: img for img in structured.get("images", [])}
    text_map = {t["id"]: t for t in structured.get("texts", [])}
    table_map = {t["id"]: t for t in structured.get("tables", [])}

    for page_layout in layout.get("pages", []):
        slide_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(slide_layout)

        for section in page_layout.get("sections", []):
            # 添加文字
            if section.get("text_id") in text_map:
                text_data = text_map[section["text_id"]]
                left, top, width, height = _get_section_position(section, prs)
                txBox = slide.shapes.add_textbox(
                    Inches(left), Inches(top), Inches(width), Inches(height)
                )
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = text_data["content"]
                p.font.size = Pt(14)

            # 添加图片
            for img_id in section.get("image_ids", []):
                if img_id in image_map:
                    img_data = image_map[img_id]
                    img_path = img_data.get("path", "")
                    if img_path and Path(img_path).exists():
                        left, top, width, height = _get_section_position(section, prs)
                        slide.shapes.add_picture(
                            img_path,
                            Inches(left), Inches(top),
                            Inches(min(width, 5)),
                        )

            # 添加表格
            for tbl_id in section.get("table_ids", []):
                if tbl_id in table_map:
                    tbl_data = table_map[tbl_id]
                    rows = len(tbl_data.get("data", []))
                    cols = len(tbl_data.get("data", [[]])[0]) if rows > 0 else 0
                    if rows > 0 and cols > 0:
                        left, top, width, height = _get_section_position(section, prs)
                        table_shape = slide.shapes.add_table(
                            rows, cols,
                            Inches(left), Inches(top),
                            Inches(width), Inches(height),
                        )
                        for r, row_data in enumerate(tbl_data["data"]):
                            for c, cell_text in enumerate(row_data):
                                table_shape.table.cell(r, c).text = str(cell_text)

    output_path = Path(settings.OUTPUT_DIR) / req.task_id / "output.pptx"
    prs.save(str(output_path))

    return GenerateResult(
        task_id=req.task_id,
        output_path=str(output_path),
        format="pptx",
        pages=len(prs.slides),
    )


def _build_typst_document(layout: dict, structured: dict) -> str:
    """将排版指令转换为 Typst 标记语言"""
    parts = []

    # 文档基础设置
    parts.append("""#set document(title: "杂志级文档", author: "智能排版系统")
#set page(paper: "a4", margin: (top: 1cm, bottom: 1cm, left: 1.5cm, right: 1.5cm))
#set text(font: ("Times New Roman", "SimSun"), size: 11pt, lang: "zh")
#set par(leading: 0.8em, justify: true)
""")

    image_map = {img["id"]: img for img in structured.get("images", [])}
    text_map = {t["id"]: t for t in structured.get("texts", [])}

    for page_layout in layout.get("pages", []):
        for section in page_layout.get("sections", []):
            # 文字
            if section.get("text_id") in text_map:
                text_data = text_map[section["text_id"]]
                content = text_data["content"]
                level = text_data.get("level", 0)
                if level == 1:
                    parts.append(f"= {content}\n")
                elif level == 2:
                    parts.append(f"== {content}\n")
                else:
                    parts.append(f"{content}\n\n")

            # 图片
            for img_id in section.get("image_ids", []):
                if img_id in image_map:
                    img_path = image_map[img_id].get("path", "")
                    if img_path and Path(img_path).exists():
                        parts.append(f'#image("{img_path}", width: 80%)\n\n')

        parts.append("#pagebreak()\n")

    return "".join(parts)


def _get_section_position(section: dict, prs) -> tuple:
    """根据布局位置计算具体坐标"""
    position = section.get("position", "top")
    slide_w = 13.333
    slide_h = 7.5

    positions = {
        "top": (1, 0.5, slide_w - 2, 2),
        "center": (1, 2.5, slide_w - 2, 3),
        "bottom": (1, 5.5, slide_w - 2, 1.5),
        "left": (0.5, 1, 5, 5.5),
        "right": (6.5, 1, 6, 5.5),
    }
    return positions.get(position, positions["top"])
