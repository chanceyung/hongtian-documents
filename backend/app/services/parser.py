"""文档解析引擎 - 本地运行，核心：图片-文字关联"""
from fastapi import APIRouter, UploadFile, File, HTTPException, BackgroundTasks
from pydantic import BaseModel
from typing import Optional
import uuid
import os
import json
from pathlib import Path

from app.core.config import settings

router = APIRouter(prefix="/parse", tags=["文档解析"])


class ParseTask(BaseModel):
    task_id: str
    status: str
    progress: float
    elements_count: int = 0
    images_count: int = 0
    tables_count: int = 0
    error: Optional[str] = None


# 任务状态存储（生产环境用 Redis）
_tasks: dict[str, ParseTask] = {}


@router.post("/upload", response_model=ParseTask)
async def upload_and_parse(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
):
    """上传文件并启动解析任务"""
    # 验证文件类型
    allowed_types = {
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/pdf",
        "text/markdown",
        "text/plain",
    }
    allowed_exts = {".pptx", ".docx", ".xlsx", ".pdf", ".md", ".txt"}

    ext = Path(file.filename or "").suffix.lower()
    if ext not in allowed_exts:
        raise HTTPException(400, f"不支持的文件类型: {ext}")

    task_id = str(uuid.uuid4())
    upload_dir = Path(settings.UPLOAD_DIR) / task_id
    upload_dir.mkdir(parents=True, exist_ok=True)

    file_path = upload_dir / file.filename
    with open(file_path, "wb") as f:
        content = await file.read()
        f.write(content)

    task = ParseTask(task_id=task_id, status="pending", progress=0)
    _tasks[task_id] = task

    background_tasks.add_task(_run_parse, task_id, file_path, ext)

    return task


async def _run_parse(task_id: str, file_path: Path, ext: str):
    """后台解析任务"""
    task = _tasks[task_id]
    try:
        task.status = "parsing"
        task.progress = 0.1

        output_dir = Path(settings.OUTPUT_DIR) / task_id
        output_dir.mkdir(parents=True, exist_ok=True)

        assets_dir = Path(settings.ASSETS_DIR) / task_id
        assets_dir.mkdir(parents=True, exist_ok=True)

        # 根据文件类型选择解析策略
        if ext == ".pdf":
            result = await _parse_pdf(file_path, assets_dir, task)
        elif ext == ".pptx":
            result = await _parse_pptx(file_path, assets_dir, task)
        elif ext == ".docx":
            result = await _parse_docx(file_path, assets_dir, task)
        elif ext == ".xlsx":
            result = await _parse_xlsx(file_path, assets_dir, task)
        else:
            result = await _parse_markdown(file_path, assets_dir, task)

        # 构建图片-文字关联关系
        task.progress = 0.8
        task.status = "linking"
        linkage = _build_content_asset_linkage(result)

        # 保存结构化结果
        structured_output = {
            "task_id": task_id,
            "source_file": str(file_path),
            "elements": result,
            "linkage": linkage,
        }

        output_file = output_dir / "structured.json"
        with open(output_file, "w", encoding="utf-8") as f:
            json.dump(structured_output, f, ensure_ascii=False, indent=2)

        task.status = "completed"
        task.progress = 1.0
        task.elements_count = len(result.get("texts", []))
        task.images_count = len(result.get("images", []))
        task.tables_count = len(result.get("tables", []))

    except Exception as e:
        task.status = "failed"
        task.error = str(e)


async def _parse_pdf(file_path: Path, assets_dir: Path, task: ParseTask) -> dict:
    """PDF 解析：PyMuPDF + PaddleOCR"""
    import fitz  # PyMuPDF

    doc = fitz.open(str(file_path))
    texts, images, tables = [], [], []

    for page_num in range(len(doc)):
        page = doc[page_num]
        task.progress = 0.1 + 0.6 * (page_num / len(doc))

        # 提取文字块及坐标
        blocks = page.get_text("dict")["blocks"]
        for block in blocks:
            if block["type"] == 0:  # 文字块
                for line in block.get("lines", []):
                    text = "".join([span["text"] for span in line.get("spans", [])])
                    if text.strip():
                        bbox = block["bbox"]
                        texts.append({
                            "id": f"text_{page_num}_{len(texts)}",
                            "page": page_num,
                            "content": text.strip(),
                            "bbox": list(bbox),
                            "font_info": _extract_font_info(block),
                        })

        # 提取图片及坐标
        image_list = page.get_images(full=True)
        for img_idx, img_info in enumerate(image_list):
            xref = img_info[0]
            base_image = doc.extract_image(xref)
            if base_image:
                img_path = assets_dir / f"page{page_num}_img{img_idx}.{base_image['ext']}"
                with open(img_path, "wb") as f:
                    f.write(base_image["image"])

                # 获取图片在页面中的位置
                img_rects = page.get_image_rects(xref)
                bbox = list(img_rects[0]) if img_rects else [0, 0, 0, 0]

                images.append({
                    "id": f"img_{page_num}_{img_idx}",
                    "page": page_num,
                    "path": str(img_path),
                    "bbox": bbox,
                    "width": base_image.get("width", 0),
                    "height": base_image.get("height", 0),
                })

    doc.close()
    return {"texts": texts, "images": images, "tables": tables}


async def _parse_pptx(file_path: Path, assets_dir: Path, task: ParseTask) -> dict:
    """PPTX 解析：python-pptx"""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation(str(file_path))
    texts, images, tables = [], []

    for slide_num, slide in enumerate(prs.slides):
        task.progress = 0.1 + 0.6 * (slide_num / len(prs.slides))

        for shape_idx, shape in enumerate(slide.shapes):
            # 文字内容
            if shape.has_text_frame:
                for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                    text = paragraph.text.strip()
                    if text:
                        texts.append({
                            "id": f"text_s{slide_num}_{shape_idx}_{para_idx}",
                            "page": slide_num,
                            "content": text,
                            "bbox": _pptx_shape_bbox(shape),
                            "level": paragraph.level,
                        })

            # 图片
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                img_bytes = shape.image.blob
                ext = shape.image.content_type.split("/")[-1]
                img_path = assets_dir / f"slide{slide_num}_shape{shape_idx}.{ext}"
                with open(img_path, "wb") as f:
                    f.write(img_bytes)

                images.append({
                    "id": f"img_s{slide_num}_{shape_idx}",
                    "page": slide_num,
                    "path": str(img_path),
                    "bbox": _pptx_shape_bbox(shape),
                })

            # 表格
            if shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                tables.append({
                    "id": f"table_s{slide_num}_{shape_idx}",
                    "page": slide_num,
                    "data": table_data,
                    "bbox": _pptx_shape_bbox(shape),
                })

    return {"texts": texts, "images": images, "tables": tables}


async def _parse_docx(file_path: Path, assets_dir: Path, task: ParseTask) -> dict:
    """Word 解析：python-docx"""
    from docx import Document

    doc = Document(str(file_path))
    texts, images, tables = [], []

    for para_idx, paragraph in enumerate(doc.paragraphs):
        task.progress = 0.1 + 0.4 * (para_idx / max(len(doc.paragraphs), 1))
        if paragraph.text.strip():
            texts.append({
                "id": f"text_p{para_idx}",
                "page": 0,
                "content": paragraph.text.strip(),
                "style": paragraph.style.name if paragraph.style else "Normal",
                "level": _docx_heading_level(paragraph),
            })

    # 提取图片
    for rel_idx, rel in enumerate(doc.part.rels.values()):
        if "image" in rel.reltype:
            img_data = rel.target_part.blob
            ext = rel.target_part.content_type.split("/")[-1]
            img_path = assets_dir / f"docx_img_{rel_idx}.{ext}"
            with open(img_path, "wb") as f:
                f.write(img_data)
            images.append({
                "id": f"img_{rel_idx}",
                "page": 0,
                "path": str(img_path),
            })

    # 提取表格
    for table_idx, table in enumerate(doc.tables):
        table_data = []
        for row in table.rows:
            row_data = [cell.text for cell in row.cells]
            table_data.append(row_data)
        tables.append({
            "id": f"table_{table_idx}",
            "page": 0,
            "data": table_data,
        })

    return {"texts": texts, "images": images, "tables": tables}


async def _parse_xlsx(file_path: Path, assets_dir: Path, task: ParseTask) -> dict:
    """Excel 解析：pandas + openpyxl"""
    import pandas as pd

    xls = pd.ExcelFile(str(file_path))
    texts, images, tables = [], []

    for sheet_name in xls.sheet_names:
        df = pd.read_excel(xls, sheet_name=sheet_name)
        tables.append({
            "id": f"table_{sheet_name}",
            "page": 0,
            "sheet": sheet_name,
            "data": df.to_dict(orient="records"),
            "columns": list(df.columns),
        })

    return {"texts": texts, "images": images, "tables": tables}


async def _parse_markdown(file_path: Path, assets_dir: Path, task: ParseTask) -> dict:
    """Markdown 解析"""
    import re

    with open(file_path, "r", encoding="utf-8") as f:
        content = f.read()

    texts = []
    for line_idx, line in enumerate(content.split("\n")):
        if line.strip():
            level = 0
            if line.startswith("#"):
                level = len(line.split(" ")[0])
            texts.append({
                "id": f"text_l{line_idx}",
                "page": 0,
                "content": line.strip(),
                "level": level,
            })

    # 提取 Markdown 中的图片引用
    images = []
    img_pattern = re.compile(r"!\[([^\]]*)\]\(([^)]+)\)")
    for match in img_pattern.finditer(content):
        images.append({
            "id": f"img_md_{len(images)}",
            "page": 0,
            "alt_text": match.group(1),
            "original_path": match.group(2),
        })

    return {"texts": texts, "images": images, "tables": []}


def _build_content_asset_linkage(parsed: dict) -> dict:
    """
    核心：构建内容-素材关联关系
    策略：空间位置 -> 语义相似度 -> 文档结构 -> 人工校验
    """
    texts = parsed.get("texts", [])
    images = parsed.get("images", [])
    tables = parsed.get("tables", [])

    linkage = {"text_image": [], "text_table": [], "low_confidence": []}

    # 策略 1：同页空间位置关联
    for text in texts:
        text_page = text.get("page", -1)
        text_bbox = text.get("bbox", [])

        for img in images:
            if img.get("page") != text_page:
                continue
            img_bbox = img.get("bbox", [])
            if not text_bbox or not img_bbox:
                continue

            distance = _bbox_distance(text_bbox, img_bbox)
            if distance < 200:  # 像素距离阈值
                confidence = max(0, 1 - distance / 200)
                linkage["text_image"].append({
                    "text_id": text["id"],
                    "image_id": img["id"],
                    "strategy": "spatial",
                    "confidence": round(confidence, 2),
                    "distance": round(distance, 1),
                })

    # 策略 2：文档结构关联（标题、图注等）
    for text in texts:
        content = text.get("content", "")
        level = text.get("level", -1)
        if level == 0 or any(kw in content for kw in ["图", "注", "见图", "如图"]):
            for img in images:
                if img.get("page") == text.get("page"):
                    linkage["text_image"].append({
                        "text_id": text["id"],
                        "image_id": img["id"],
                        "strategy": "structural",
                        "confidence": 0.7,
                    })

    # 策略 3：表格-文字关联
    for text in texts:
        for table in tables:
            if text.get("page") == table.get("page"):
                content = text.get("content", "")
                if any(kw in content for kw in ["表", "数据", "如下", "见表"]):
                    linkage["text_table"].append({
                        "text_id": text["id"],
                        "table_id": table["id"],
                        "strategy": "structural",
                        "confidence": 0.6,
                    })

    # 标记低置信度关联
    for link in linkage["text_image"]:
        if link["confidence"] < 0.5:
            linkage["low_confidence"].append({
                "type": "text_image",
                "text_id": link["text_id"],
                "image_id": link["image_id"],
                "confidence": link["confidence"],
                "reason": "需要人工确认关联关系",
            })

    return linkage


def _bbox_distance(bbox1: list, bbox2: list) -> float:
    """计算两个边界框之间的距离"""
    if len(bbox1) < 4 or len(bbox2) < 4:
        return float("inf")

    cx1 = (bbox1[0] + bbox1[2]) / 2
    cy1 = (bbox1[1] + bbox1[3]) / 2
    cx2 = (bbox2[0] + bbox2[2]) / 2
    cy2 = (bbox2[1] + bbox2[3]) / 2

    return ((cx1 - cx2) ** 2 + (cy1 - cy2) ** 2) ** 0.5


def _extract_font_info(block: dict) -> dict:
    """提取字体信息"""
    fonts = set()
    sizes = set()
    for line in block.get("lines", []):
        for span in line.get("spans", []):
            if span.get("font"):
                fonts.add(span["font"])
            if span.get("size"):
                sizes.add(round(span["size"], 1))
    return {"fonts": list(fonts), "sizes": list(sizes)}


def _pptx_shape_bbox(shape) -> list:
    """PPTX shape 转为 bbox"""
    try:
        return [
            shape.left, shape.top,
            shape.left + shape.width,
            shape.top + shape.height,
        ]
    except Exception:
        return [0, 0, 0, 0]


def _docx_heading_level(paragraph) -> int:
    """Word 段落标题级别"""
    style_name = (paragraph.style.name or "").lower()
    if "heading" in style_name:
        for i in range(1, 7):
            if str(i) in style_name:
                return i
    return 0


@router.get("/status/{task_id}", response_model=ParseTask)
async def get_parse_status(task_id: str):
    """查询解析任务状态"""
    task = _tasks.get(task_id)
    if not task:
        raise HTTPException(404, "任务不存在")
    return task


@router.get("/result/{task_id}")
async def get_parse_result(task_id: str):
    """获取解析结果"""
    result_file = Path(settings.OUTPUT_DIR) / task_id / "structured.json"
    if not result_file.exists():
        raise HTTPException(404, "解析结果不存在")

    with open(result_file, "r", encoding="utf-8") as f:
        return json.load(f)
