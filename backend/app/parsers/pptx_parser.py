"""PPTX 解析 — python-pptx"""
import hashlib
from pathlib import Path

from pptx import Presentation
from pptx.shapes.graphfrm import GraphicFrame
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.models.unified_document import (
    UnifiedDocument,
    TextElement,
    ImageElement,
    TableElement,
    BoundingBox,
    PageLayout,
)


class PptxParser:

    async def parse(self, path: Path, session_id: str) -> UnifiedDocument:
        assets_dir = path.parent / "assets" / session_id
        assets_dir.mkdir(parents=True, exist_ok=True)

        prs = Presentation(str(path))
        warnings: list[str] = []
        texts: list[TextElement] = []
        images: list[ImageElement] = []
        tables: list[TableElement] = []
        page_layouts: list[PageLayout] = []

        slide_width = prs.slide_width
        slide_height = prs.slide_height

        for slide_idx, slide in enumerate(prs.slides):
            layout_info = self._extract_layout(slide, slide_idx, slide_width, slide_height)

            for shape in slide.shapes:
                try:
                    bbox = BoundingBox(
                        left=shape.left, top=shape.top,
                        width=shape.width, height=shape.height,
                    )

                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        images.extend(self._parse_image(shape, slide_idx, bbox, assets_dir))
                    elif hasattr(shape, "has_table") and shape.has_table:
                        tables.extend(self._parse_table(shape, slide_idx, bbox))
                    elif hasattr(shape, "has_chart") and shape.has_chart:
                        tables.extend(self._parse_chart(shape, slide_idx, bbox))
                    elif hasattr(shape, "has_text_frame") and shape.has_text_frame:
                        texts.extend(self._parse_text(shape, slide_idx, bbox))
                except Exception as e:
                    warnings.append(f"Slide {slide_idx} Shape {shape.shape_id}: {e}")

            # Build visual hierarchy from sorted shapes (by area, descending)
            sorted_ids = sorted(
                [s for s in slide.shapes],
                key=lambda s: s.width * s.height,
                reverse=True,
            )
            layout_info.visual_hierarchy = [f"s{slide_idx}_sh{s.shape_id}" for s in sorted_ids[:10]]
            page_layouts.append(layout_info)

        return UnifiedDocument(
            source_file=str(path),
            source_format="pptx",
            parse_method="python-pptx",
            total_pages=len(prs.slides),
            texts=texts,
            images=images,
            tables=tables,
            parse_warnings=warnings,
            page_layouts=page_layouts,
        )

    def _parse_text(self, shape, slide_idx: int, bbox: BoundingBox) -> list[TextElement]:
        elements = []
        for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
            text = paragraph.text.strip()
            if not text:
                continue

            level = 0
            if paragraph.level is not None and paragraph.level > 0:
                level = paragraph.level

            original_font = None
            original_size = None
            original_color = None
            if paragraph.runs:
                run = paragraph.runs[0]
                if run.font.name:
                    original_font = run.font.name
                if run.font.size:
                    original_size = run.font.size.pt
                if run.font.color and run.font.color.rgb:
                    original_color = f"#{run.font.color.rgb}"

            elements.append(TextElement(
                id=f"s{slide_idx}_sh{shape.shape_id}_p{para_idx}",
                content=text,
                page=slide_idx,
                bbox=bbox,
                level=level,
                fingerprint=hashlib.md5(text.encode()).hexdigest(),
                original_font=original_font,
                original_size=original_size,
                original_color=original_color,
            ))
        return elements

    def _parse_image(self, shape, slide_idx: int, bbox: BoundingBox, assets_dir: Path) -> list[ImageElement]:
        try:
            blob = shape.image.blob
            ext = shape.image.ext
            img_hash = hashlib.md5(blob).hexdigest()[:12]
            filename = f"s{slide_idx}_sh{shape.shape_id}_{img_hash}.{ext}"
            img_path = assets_dir / filename
            img_path.write_bytes(blob)

            return [ImageElement(
                id=f"s{slide_idx}_sh{shape.shape_id}",
                local_path=str(img_path),
                page=slide_idx,
                bbox=bbox,
                hash=img_hash,
                alt_text=shape.name or "",
            )]
        except Exception:
            return []

    def _parse_table(self, shape: GraphicFrame, slide_idx: int, bbox: BoundingBox) -> list[TableElement]:
        table = shape.table
        headers: list[str] = []
        data: list[list[str]] = []

        for row_idx, row in enumerate(table.rows):
            row_data = [cell.text.strip() for cell in row.cells]
            if row_idx == 0:
                headers = row_data
            else:
                data.append(row_data)

        return [TableElement(
            id=f"s{slide_idx}_sh{shape.shape_id}_tbl",
            page=slide_idx,
            bbox=bbox,
            data=data,
            headers=headers,
        )]

    def _parse_chart(self, shape: GraphicFrame, slide_idx: int, bbox: BoundingBox) -> list[TableElement]:
        try:
            chart = shape.chart
            data: list[list[str]] = []
            headers = ["Series"]

            series_list = list(chart.series)
            if series_list:
                categories = [str(c) for c in series_list[0].categories]
                headers.extend(categories)

            for series in series_list:
                row_data = [str(series.name)]
                for point in series.points:
                    row_data.append(str(point.value) if point.value is not None else "")
                data.append(row_data)

            return [TableElement(
                id=f"s{slide_idx}_sh{shape.shape_id}_chart",
                page=slide_idx,
                bbox=bbox,
                data=data,
                headers=headers,
                is_chart=True,
            )]
        except Exception:
            return []

    def _extract_layout(self, slide, slide_idx: int, slide_width: int, slide_height: int) -> PageLayout:
        layout_type = "unknown"
        try:
            layout_name = slide.slide_layout.name.lower() if slide.slide_layout else ""
            if "title" in layout_name or "cover" in layout_name:
                layout_type = "cover"
            elif "two" in layout_name or "column" in layout_name:
                layout_type = "two_column"
            elif "blank" in layout_name:
                layout_type = "full_text"
            elif "content" in layout_name:
                layout_type = "text_image"
        except Exception:
            pass

        shape_area = sum(s.width * s.height for s in slide.shapes) if slide.shapes else 0
        total_area = slide_width * slide_height if slide_width and slide_height else 1
        whitespace_ratio = max(0.0, 1.0 - shape_area / total_area) if total_area > 0 else 0.0

        return PageLayout(
            page_number=slide_idx,
            layout_type=layout_type,
            whitespace_ratio=round(whitespace_ratio, 3),
            original_structure={"layout_name": getattr(slide.slide_layout, "name", "")},
        )
