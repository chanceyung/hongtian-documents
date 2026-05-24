"""验证门控系统 — 在关键节点验证处理质量"""
from __future__ import annotations

from pathlib import Path
from typing import Literal

from pydantic import BaseModel

from app.core.logging import get_logger
from app.models.edit_actions import MagazineEditPlan
from app.models.unified_document import ContentFingerprint, UnifiedDocument

# Import these at module level for easier mocking in tests
try:
    import pypdf
except ImportError:
    pypdf = None  # type: ignore

try:
    from pptx import Presentation  # type: ignore
except ImportError:
    Presentation = None  # type: ignore

logger = get_logger(__name__)


class GateResult(BaseModel):
    """单个验证门的结果"""
    gate_name: str
    passed: bool
    score: float
    issues: list[str]
    warnings: list[str] = []


class ValidationGates:
    """三道验证门控系统"""

    def __init__(self, text_threshold: float = 0.98) -> None:
        self.text_threshold = text_threshold
        self.logger = logger

    async def gate_1_parse_completeness(self, doc: UnifiedDocument) -> GateResult:
        """
        Gate 1 (after Parser): 检查解析完整性

        验证项:
        - 文本提取率 >= 98%
        - 所有提取的图片在磁盘上存在
        - 所有表格有数据
        """
        issues: list[str] = []
        warnings: list[str] = []
        score = 1.0

        # 1. 文本提取率检查
        text_pages = len({t.page for t in doc.texts})
        total_pages = max(doc.total_pages, 1)
        text_rate = text_pages / total_pages if total_pages > 0 else 0

        if text_rate < self.text_threshold:
            issue = f"文本提取率不足: {text_rate:.2%} < {self.text_threshold:.2%} ({text_pages}/{total_pages} 页)"
            issues.append(issue)
            score -= (self.text_threshold - text_rate) * 2
            self.logger.warning(issue)

        # 2. 图片文件存在性检查
        missing_images: list[str] = []
        for img in doc.images:
            if not Path(img.local_path).exists():
                missing_images.append(img.id)
                issues.append(f"图片文件不存在: {img.local_path}")

        if missing_images:
            score -= len(missing_images) * 0.1
            self.logger.warning(f"发现 {len(missing_images)} 个缺失图片文件")

        # 3. 表格数据完整性检查
        empty_tables: list[str] = []
        for table in doc.tables:
            if not table.data or len(table.data) == 0:
                empty_tables.append(table.id)
                issues.append(f"表格无数据: {table.id}")

        if empty_tables:
            score -= len(empty_tables) * 0.1
            self.logger.warning(f"发现 {len(empty_tables)} 个空表格")

        # 4. 警告：页数异常
        if doc.total_pages == 0:
            warnings.append("文档总页数为 0，可能解析失败")

        # 5. 警告：无内容
        if len(doc.texts) == 0 and len(doc.images) == 0 and len(doc.tables) == 0:
            warnings.append("文档未提取到任何内容")

        score = max(0.0, min(1.0, score))
        passed = len(issues) == 0

        self.logger.info(
            "Gate 1 解析完整性检查完成",
            passed=passed,
            score=score,
            issues_count=len(issues),
            warnings_count=len(warnings),
        )

        return GateResult(
            gate_name="gate_1_parse_completeness",
            passed=passed,
            score=score,
            issues=issues,
            warnings=warnings,
        )

    async def gate_2_content_understanding(
        self,
        doc: UnifiedDocument,
        edit_plan: MagazineEditPlan,
        fingerprint: ContentFingerprint,
    ) -> GateResult:
        """
        Gate 2 (after Analyzer): 检查内容理解

        验证项:
        - 内容组覆盖所有原始文本元素
        - 内容组覆盖所有原始图片元素
        - 无孤立元素（未被任何编辑计划引用）
        """
        issues: list[str] = []
        warnings: list[str] = []
        score = 1.0

        # 收集所有被引用的源 ID
        referenced_text_ids: set[str] = set()
        referenced_image_ids: set[str] = set()

        for page_plan in edit_plan.pages:
            for action in page_plan.actions:
                if action.source_id in fingerprint.text_fingerprints:
                    referenced_text_ids.add(action.source_id)
                elif action.source_id in fingerprint.image_hashes:
                    referenced_image_ids.add(action.source_id)

        # 1. 文本元素覆盖率检查
        all_text_ids = set(fingerprint.text_fingerprints.keys())
        orphaned_texts = all_text_ids - referenced_text_ids

        if orphaned_texts:
            orphan_rate = len(orphaned_texts) / len(all_text_ids) if all_text_ids else 0
            if orphan_rate > 0.1:  # 超过 10% 未覆盖
                issue = f"{len(orphaned_texts)} 个文本元素未被覆盖 ({orphan_rate:.1%})"
                issues.append(issue)
                score -= orphan_rate
                self.logger.warning(issue)
            else:
                warnings.append(f"少量文本元素未覆盖: {len(orphaned_texts)} 个")

        # 2. 图片元素覆盖率检查
        all_image_ids = set(fingerprint.image_hashes.keys())
        orphaned_images = all_image_ids - referenced_image_ids

        if orphaned_images:
            orphan_rate = len(orphaned_images) / len(all_image_ids) if all_image_ids else 0
            if orphan_rate > 0.1:
                issue = f"{len(orphaned_images)} 个图片元素未被覆盖 ({orphan_rate:.1%})"
                issues.append(issue)
                score -= orphan_rate
                self.logger.warning(issue)
            else:
                warnings.append(f"少量图片元素未覆盖: {len(orphaned_images)} 个")

        # 3. 检查页数匹配
        expected_pages = doc.total_pages
        actual_pages = len(edit_plan.pages)

        if actual_pages != expected_pages:
            issue = f"页数不匹配: 编辑计划 {actual_pages} 页 vs 原文档 {expected_pages} 页"
            issues.append(issue)
            score -= abs(actual_pages - expected_pages) * 0.1
            self.logger.warning(issue)

        # 4. 警告：空编辑计划
        if len(edit_plan.pages) == 0:
            warnings.append("编辑计划为空，未生成任何页面")

        # 5. 警告：页面无动作
        empty_pages = [p.page_number for p in edit_plan.pages if not p.actions]
        if empty_pages:
            warnings.append(f"页面无编辑动作: {empty_pages}")

        score = max(0.0, min(1.0, score))
        passed = len(issues) == 0

        self.logger.info(
            "Gate 2 内容理解检查完成",
            passed=passed,
            score=score,
            issues_count=len(issues),
            warnings_count=len(warnings),
            referenced_texts=len(referenced_text_ids),
            orphaned_texts=len(orphaned_texts),
            referenced_images=len(referenced_image_ids),
            orphaned_images=len(orphaned_images),
        )

        return GateResult(
            gate_name="gate_2_content_understanding",
            passed=passed,
            score=score,
            issues=issues,
            warnings=warnings,
        )

    async def gate_3_render_quality(
        self,
        output_path: Path,
        output_format: Literal["pdf", "pptx"],
    ) -> GateResult:
        """
        Gate 3 (after Renderer): 检查渲染质量

        验证项:
        - 输出文件存在且大小 > 0
        - PDF 可被 PyPDF2 打开
        - PPTX 可被 python-pptx 打开
        - 基本文件完整性检查
        """
        issues: list[str] = []
        warnings: list[str] = []
        score = 1.0

        # 1. 文件存在性和大小检查
        if not output_path.exists():
            issues.append(f"输出文件不存在: {output_path}")
            self.logger.error(f"输出文件不存在: {output_path}")
            return GateResult(
                gate_name="gate_3_render_quality",
                passed=False,
                score=0.0,
                issues=issues,
                warnings=warnings,
            )

        file_size = output_path.stat().st_size
        if file_size == 0:
            issues.append(f"输出文件大小为 0: {output_path}")
            self.logger.error(f"输出文件大小为 0: {output_path}")
            return GateResult(
                gate_name="gate_3_render_quality",
                passed=False,
                score=0.0,
                issues=issues,
                warnings=warnings,
            )

        # 2. 格式特定的完整性检查
        if output_format == "pdf":
            if pypdf is None:
                issues.append("pypdf 模块未安装，无法验证 PDF")
                score -= 1.0
            else:
                try:
                    with output_path.open("rb") as f:
                        reader = pypdf.PdfReader(f)
                        if len(reader.pages) == 0:
                            issues.append("PDF 文件无页面")
                            score -= 0.5
                        else:
                            self.logger.info(f"PDF 验证通过: {len(reader.pages)} 页")
                except Exception as e:
                    issues.append(f"PDF 文件损坏或无法打开: {e}")
                    score -= 1.0
                    self.logger.error(f"PDF 验证失败: {e}")

        elif output_format == "pptx":
            if Presentation is None:
                issues.append("python-pptx 模块未安装，无法验证 PPTX")
                score -= 1.0
            else:
                try:
                    prs = Presentation(output_path)
                    if len(prs.slides) == 0:
                        issues.append("PPTX 文件无幻灯片")
                        score -= 0.5
                    else:
                        self.logger.info(f"PPTX 验证通过: {len(prs.slides)} 张幻灯片")
                except Exception as e:
                    issues.append(f"PPTX 文件损坏或无法打开: {e}")
                    score -= 1.0
                    self.logger.error(f"PPTX 验证失败: {e}")
        else:
            warnings.append(f"未知的输出格式: {output_format}")

        # 3. 警告：文件过小
        if file_size < 1024:  # 小于 1KB
            warnings.append(f"输出文件过小: {file_size} 字节")

        # 4. 警告：文件扩展名不匹配
        expected_ext = f".{output_format}"
        if output_path.suffix.lower() != expected_ext:
            warnings.append(f"文件扩展名不匹配: {output_path.suffix} vs {expected_ext}")

        score = max(0.0, min(1.0, score))
        passed = len(issues) == 0

        self.logger.info(
            "Gate 3 渲染质量检查完成",
            passed=passed,
            score=score,
            issues_count=len(issues),
            warnings_count=len(warnings),
            file_size=file_size,
            output_path=str(output_path),
        )

        return GateResult(
            gate_name="gate_3_render_quality",
            passed=passed,
            score=score,
            issues=issues,
            warnings=warnings,
        )

    async def run_all_gates(
        self,
        doc: UnifiedDocument,
        edit_plan: MagazineEditPlan,
        fingerprint: ContentFingerprint,
        output_path: Path,
        output_format: Literal["pdf", "pptx"],
    ) -> dict[str, GateResult]:
        """运行所有三道验证门"""
        self.logger.info("开始运行所有验证门")

        gate1_result = await self.gate_1_parse_completeness(doc)
        gate2_result = await self.gate_2_content_understanding(doc, edit_plan, fingerprint)
        gate3_result = await self.gate_3_render_quality(output_path, output_format)

        results = {
            "gate_1": gate1_result,
            "gate_2": gate2_result,
            "gate_3": gate3_result,
        }

        all_passed = all(r.passed for r in results.values())
        avg_score = sum(r.score for r in results.values()) / len(results)

        self.logger.info(
            "所有验证门检查完成",
            all_passed=all_passed,
            avg_score=avg_score,
            gate1_passed=gate1_result.passed,
            gate2_passed=gate2_result.passed,
            gate3_passed=gate3_result.passed,
        )

        return results